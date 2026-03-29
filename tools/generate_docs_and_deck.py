import os
import re
import textwrap
from datetime import datetime
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = REPO_ROOT / "docs"
LIB_DIR = REPO_ROOT / "lib"
FUNCTIONS_INDEX = REPO_ROOT / "functions" / "src" / "index.ts"
PUBSPEC = REPO_ROOT / "pubspec.yaml"
FUNCTIONS_PACKAGE_JSON = REPO_ROOT / "functions" / "package.json"


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


def _md_h1(title: str) -> str:
    return f"# {title}\n\n"


def _md_h2(title: str) -> str:
    return f"## {title}\n\n"


def _md_h3(title: str) -> str:
    return f"### {title}\n\n"


def _bullet(items: list[str]) -> str:
    return "".join([f"- {i}\n" for i in items]) + "\n"


def _codeblock(code: str, lang: str = "") -> str:
    fence = "```"
    return f"{fence}{lang}\n{code.rstrip()}\n{fence}\n\n"


def load_pubspec_dependencies() -> dict:
    data = yaml.safe_load(_read_text(PUBSPEC))
    deps = data.get("dependencies", {}) or {}
    dev_deps = data.get("dev_dependencies", {}) or {}
    return {"dependencies": deps, "dev_dependencies": dev_deps}


def load_functions_dependencies() -> dict:
    if not FUNCTIONS_PACKAGE_JSON.exists():
        return {}
    data = __import__("json").loads(_read_text(FUNCTIONS_PACKAGE_JSON))
    return {
        "engines": data.get("engines", {}),
        "dependencies": data.get("dependencies", {}),
        "devDependencies": data.get("devDependencies", {}),
    }


def list_dart_files(folder: Path) -> list[Path]:
    return sorted([p for p in folder.rglob("*.dart") if p.is_file()])


def extract_classes(dart_source: str) -> list[str]:
    # Basic extractor for "class X" declarations
    return re.findall(r"^\s*class\s+([A-Za-z0-9_]+)\b", dart_source, flags=re.M)


def extract_firestore_collections() -> dict:
    pattern = re.compile(r"\.collection(?:Group)?\('([^']+)'\)")
    results: dict[str, set[str]] = {"collection": set(), "collectionGroup": set()}

    for path in [*list_dart_files(LIB_DIR), FUNCTIONS_INDEX]:
        if not path.exists():
            continue
        txt = _read_text(path)
        for m in re.finditer(r"\.collection\('([^']+)'\)", txt):
            results["collection"].add(m.group(1))
        for m in re.finditer(r"\.collectionGroup\('([^']+)'\)", txt):
            results["collectionGroup"].add(m.group(1))

    return {k: sorted(v) for k, v in results.items()}


def extract_functions_exports() -> list[str]:
    if not FUNCTIONS_INDEX.exists():
        return []
    txt = _read_text(FUNCTIONS_INDEX)
    return re.findall(r"export\s+const\s+([A-Za-z0-9_]+)\s*=", txt)


def summarize_folder(name: str, folder: Path) -> list[dict]:
    items = []
    for path in sorted(folder.glob("*.dart")):
        src = _read_text(path)
        classes = extract_classes(src)
        items.append(
            {
                "file": str(path.relative_to(REPO_ROOT)).replace("\\", "/"),
                "classes": classes,
                "lines": src.count("\n") + 1,
            }
        )
    return items


def extract_top_level_methods(dart_source: str) -> list[str]:
    # Heuristic: captures method names that look like "Future<...> foo(" or "void foo(" etc.
    # Not a full Dart parser, but good enough for documentation inventory.
    names = set()
    for m in re.finditer(
        r"^\s*(?:Future<[^>]+>|Future|Stream<[^>]+>|Stream|void|bool|int|double|String|Map<[^>]+>|List<[^>]+>|Widget)\s+([A-Za-z0-9_]+)\s*\(",
        dart_source,
        flags=re.M,
    ):
        names.add(m.group(1))
    return sorted(names)


def build_detailed_file_guide(dart_files: list[Path], ts_files: list[Path]) -> str:
    md = _md_h1("Detailed Code Guide (Per File)")
    md += (
        "This section is designed for supervisors, auditors, and engineering reviewers. "
        "It explains the purpose of each major file and highlights the key classes and methods.\n\n"
    )

    def add_file(path: Path, kind: str) -> None:
        nonlocal md
        rel = str(path.relative_to(REPO_ROOT)).replace("\\", "/")
        src = _read_text(path)
        classes = extract_classes(src) if path.suffix == ".dart" else re.findall(r"export\s+const\s+([A-Za-z0-9_]+)\s*=", src)
        methods = extract_top_level_methods(src) if path.suffix == ".dart" else []
        md += _md_h2(rel)
        md += _bullet(
            [
                f"Type: {kind}",
                f"Lines: {src.count(chr(10)) + 1}",
                f"Classes/Exports: {', '.join(classes) if classes else '(none detected)'}",
                f"Key methods: {', '.join(methods[:30]) + ('...' if len(methods) > 30 else '') if methods else '(n/a)'}",
            ]
        )
        # Add a small excerpt (top of file) to give context without duplicating full source.
        excerpt = "\n".join(src.splitlines()[:60])
        md += _md_h3("Header Excerpt")
        md += _codeblock(excerpt, "dart" if path.suffix == ".dart" else "ts")

    # Group by folders for readability
    for folder in ["app", "providers", "screens", "services", "widgets", "models"]:
        group = [p for p in dart_files if f"lib/{folder}/" in str(p.as_posix())]
        if not group:
            continue
        md += _md_h2(f"lib/{folder}/")
        for p in group:
            add_file(p, kind="Flutter/Dart")

    if ts_files:
        md += _md_h2("functions/src/")
        for p in ts_files:
            add_file(p, kind="Firebase Cloud Functions (TypeScript)")

    return md


def build_source_appendix(dart_files: list[Path], ts_files: list[Path]) -> str:
    md = _md_h1("Appendix: Source Code (Reference)")
    md += (
        "This appendix includes the source code as a reference for deep reviews. "
        "It is intended to be printed/exported into a long-form document if needed.\n\n"
    )

    def add_source(path: Path, lang: str) -> None:
        nonlocal md
        rel = str(path.relative_to(REPO_ROOT)).replace("\\", "/")
        md += _md_h2(rel)
        md += _codeblock(_read_text(path), lang)

    for p in dart_files:
        add_source(p, "dart")
    for p in ts_files:
        add_source(p, "ts")
    return md


def build_docs() -> dict[str, Path]:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)

    stamp = datetime.now().strftime("%Y-%m-%d")
    deps = load_pubspec_dependencies()
    fn_deps = load_functions_dependencies()
    firestore = extract_firestore_collections()
    fn_exports = extract_functions_exports()

    screens = summarize_folder("screens", LIB_DIR / "screens")
    providers = summarize_folder("providers", LIB_DIR / "providers")
    services = summarize_folder("services", LIB_DIR / "services")
    widgets = summarize_folder("widgets", LIB_DIR / "widgets")
    models = summarize_folder("models", LIB_DIR / "models")

    # 00 Executive summary
    exec_md = (
        _md_h1("Cartiqo (SwiftCart) Project Summary")
        + f"_Generated on {stamp} from the current codebase._\n\n"
        + _md_h2("One-Line Description")
        + "Cartiqo is a Flutter + Firebase retail checkout and mall-operations platform: a customer-facing scan-to-bill app and a web portal for admins and mall managers to manage malls, subscriptions, managers, inventory, QR/barcode assets, and sales history.\n\n"
        + _md_h2("What Exists In This Repo")
        + _bullet(
            [
                "Flutter app (Android/iOS/Web) built with Provider state management.",
                "Web portal homepage with Admin login and Mall Manager login entry points.",
                "Firestore-backed data model for malls, managers, products, barcodes, payments, promotions, announcements, staff activity, support requests, and user profiles.",
                "CSV export utilities and PDF/print utilities for labels and receipts.",
                "Firebase Cloud Functions (Node.js/TypeScript) for OTP flows (phone OTP SMS, email OTP signup, password reset OTP) plus scheduled cleanup jobs.",
            ]
        )
        + _md_h2("Primary Roles and Portals")
        + _bullet(
            [
                "Customer/App User: signs up/logs in (email/password or OTP), selects a mall, scans items, builds cart, and checks out.",
                "Admin (Web): manages malls, subscriptions, managers, and operational dashboards; monitors support requests and announcements.",
                "Mall Manager (Web): logs in via manager ID + password; manages products, barcodes, billing settings, promotions, and sales analytics with CSV export.",
            ]
        )
    )
    p_exec = DOCS_DIR / "00_EXECUTIVE_SUMMARY.md"
    _write_text(p_exec, exec_md)

    # 01 Tools/stack
    dep_lines = []
    for k, v in deps["dependencies"].items():
        if isinstance(v, dict):
            dep_lines.append(f"{k}: (sdk)")
        else:
            dep_lines.append(f"{k}: {v}")
    dev_dep_lines = []
    for k, v in deps["dev_dependencies"].items():
        if isinstance(v, dict):
            dev_dep_lines.append(f"{k}: (sdk)")
        else:
            dev_dep_lines.append(f"{k}: {v}")

    tools_md = (
        _md_h1("Models, Tools, and Technologies")
        + f"_Generated on {stamp} from the current codebase._\n\n"
        + _md_h2("Core App Stack")
        + _bullet(
            [
                "Flutter (Dart) for Android/iOS/Web UI.",
                "Provider for state management.",
                "Firebase Core + Cloud Firestore as the primary backend database.",
                "Firebase Auth for user authentication.",
                "Firebase Hosting for web deployment.",
            ]
        )
        + _md_h2("Operational/Feature Libraries")
        + _bullet(
            [
                "mobile_scanner + vibration: barcode/QR scanning with haptics feedback.",
                "pdf + printing: generate printable PDFs (labels/receipts) and trigger print dialogs.",
                "csv + file_picker + web download helper: export operational data (sales dashboard, subscriptions, etc.) as CSV.",
                "crypto (sha256): hashing for manager/admin credential workflows stored in Firestore (as implemented in providers).",
                "qr_flutter: QR generation for mall assets and labels.",
            ]
        )
        + _md_h2("Firebase Cloud Functions Stack (functions/)")
        + _bullet(
            [
                f"Runtime: {fn_deps.get('engines', {}).get('node', '(not set)')}",
                "TypeScript build pipeline (`npm run build`).",
                "firebase-admin + firebase-functions.",
                "Twilio for SMS OTP sending (requires Twilio env vars).",
                "Resend (HTTP API) for email OTP flows (requires API key).",
            ]
        )
        + _md_h2("Flutter Dependencies (pubspec.yaml)")
        + _md_h3("Dependencies")
        + _codeblock("\n".join(dep_lines), "yaml")
        + _md_h3("Dev Dependencies")
        + _codeblock("\n".join(dev_dep_lines), "yaml")
    )
    p_tools = DOCS_DIR / "01_TOOLS_TECH_STACK.md"
    _write_text(p_tools, tools_md)

    # 02 Workflows
    workflows_md = (
        _md_h1("App and Website Workflows")
        + f"_Generated on {stamp} from the current codebase._\n\n"
        + _md_h2("Customer App Workflow (Mobile)")
        + _bullet(
            [
                "Auth: Email/password signup + email verification OR phone OTP login.",
                "Mall selection: choose a mall from Firestore (`malls` collection).",
                "Scan-to-cart: barcode scanner reads product barcode; maps barcode to product in the selected mall.",
                "Cart management: add/remove items and adjust quantities.",
                "Checkout: billing settings (GST/tax/extra charges) applied; payment record is created in Firestore under mall payments.",
                "History: user profile provides access to previous bills and payments (per-user subcollections).",
            ]
        )
        + _md_h2("Admin Web Portal Workflow")
        + _bullet(
            [
                "Admin login (web) then access Admin dashboard.",
                "Create and manage malls and subscriptions (also stored in `mall_subscriptions`).",
                "Manage mall managers (multiple managers per mall): create/update/reset credentials, view manager list for a mall.",
                "Monitor recent activity and payments (uses `collectionGroup('payments')` watch).",
                "Announcements + support requests workflows (create announcements, triage support tickets).",
                "Exporting: subscription/summary and other admin tables can be exported as CSV (where enabled).",
            ]
        )
        + _md_h2("Mall Manager Web Portal Workflow")
        + _bullet(
            [
                "Manager login using manager ID + password (lookup via `manager_index` and `malls/{mallId}/managers`).",
                "Inventory: create/update/delete products in `malls/{mallId}/products`.",
                "Barcode mapping: maintain `malls/{mallId}/barcodes` mapping for fast scanning.",
                "Billing settings: update GST/tax/extra charges stored under the mall document (`billingSettings`).",
                "Promotions: create and manage promotions under `malls/{mallId}/promotions`.",
                "Sales dashboard: filter and consolidate payments, group analysis (day/week/month/year), refresh button, and CSV export.",
                "Staff activity logging: key actions recorded under `malls/{mallId}/staff_activity`.",
            ]
        )
    )
    p_workflows = DOCS_DIR / "02_WORKFLOWS.md"
    _write_text(p_workflows, workflows_md)

    # 03 Data models
    models_md = _md_h1("Data Models (Dart)")
    models_md += f"_Generated on {stamp} from lib/models._\n\n"
    models_md += _md_h2("Model Files")
    models_md += _bullet([f"{m['file']} ({m['lines']} lines)" for m in models])
    models_md += _md_h2("Notes")
    models_md += (
        "These are the primary client-side data structures used to map Firestore documents into strongly-typed Dart objects.\n\n"
    )
    for m in models:
        models_md += _md_h3(m["file"])
        models_md += _bullet([f"Classes: {', '.join(m['classes']) if m['classes'] else '(none detected)'}"])
    p_models = DOCS_DIR / "03_DATA_MODELS.md"
    _write_text(p_models, models_md)

    # 04 Firestore schema (from usage)
    schema_md = (
        _md_h1("Firestore Collections and Schema (Inferred)")
        + f"_Generated on {stamp} by scanning Firestore calls in lib/ and functions/._\n\n"
        + _md_h2("Top-Level Collections Referenced")
        + _bullet(firestore.get("collection", []))
        + _md_h2("Collection Group Queries Referenced")
        + _bullet(firestore.get("collectionGroup", []))
        + _md_h2("Known Subcollection Patterns (Observed in Code)")
        + _bullet(
            [
                "malls/{mallId}/products",
                "malls/{mallId}/barcodes",
                "malls/{mallId}/payments",
                "malls/{mallId}/managers",
                "malls/{mallId}/promotions",
                "malls/{mallId}/staff_activity",
                "users/{uid}/bills",
                "users/{uid}/payments",
            ]
        )
        + _md_h2("OTP/Support Collections (Used by Auth + Functions)")
        + _bullet(
            [
                "phone_otp (SMS OTP documents)",
                "email_otp_signup (email signup OTP documents)",
                "password_reset_otp (password reset OTP documents)",
                "support_requests (web portal contact/demo requests)",
            ]
        )
    )
    p_schema = DOCS_DIR / "04_FIRESTORE_SCHEMA.md"
    _write_text(p_schema, schema_md)

    # 05 Code walkthrough (auto inventory)
    inv_md = (
        _md_h1("Code Walkthrough (Inventory)")
        + f"_Generated on {stamp} from the current lib/ layout._\n\n"
        + _md_h2("Screens (UI)")
        + _bullet([f"{s['file']}: {', '.join(s['classes'])}" for s in screens])
        + _md_h2("Providers (State + Firestore)")
        + _bullet([f"{p['file']}: {', '.join(p['classes'])}" for p in providers])
        + _md_h2("Services (Export/Print/Utilities)")
        + _bullet([f"{s['file']}: {', '.join(s['classes'])}" for s in services])
        + _md_h2("Widgets (Reusable UI Components)")
        + _bullet([f"{w['file']}: {', '.join(w['classes'])}" for w in widgets])
        + _md_h2("Cloud Functions Exports")
        + _bullet(fn_exports if fn_exports else ["(functions/src/index.ts not found)"])
    )
    p_inv = DOCS_DIR / "05_CODE_WALKTHROUGH.md"
    _write_text(p_inv, inv_md)

    # 06 Deployment
    deploy_md = (
        _md_h1("Deployment and Operations")
        + f"_Generated on {stamp}._\n\n"
        + _md_h2("Web Hosting (Firebase Hosting)")
        + "Build + deploy steps:\n\n"
        + _codeblock(
            "\n".join(
                [
                    "cd C:\\Users\\GS\\swiftcart_app",
                    "flutter pub get",
                    "flutter build web --release",
                    "firebase deploy --only hosting",
                ]
            ),
            "powershell",
        )
        + _md_h2("Cloud Functions (Firebase Functions)")
        + "Local build steps:\n\n"
        + _codeblock(
            "\n".join(
                [
                    "cd C:\\Users\\GS\\swiftcart_app\\functions",
                    "npm ci",
                    "npm run build",
                    "firebase deploy --only functions",
                ]
            ),
            "powershell",
        )
        + _md_h3("Important: Billing Plan Requirement")
        + "Deploying Cloud Functions can require enabling Google APIs (Cloud Build, Artifact Registry). If your Firebase project is on the free plan, deploy may be blocked until upgraded to Blaze.\n\n"
        + _md_h2("Required Secrets/Environment Variables (Functions)")
        + _bullet(
            [
                "TWILIO_ACCOUNT_SID",
                "TWILIO_AUTH_TOKEN",
                "TWILIO_PHONE_NUMBER",
                "RESEND_API_KEY (if using Resend email flow)",
                "FROM_EMAIL (optional override)",
            ]
        )
    )
    p_deploy = DOCS_DIR / "06_DEPLOYMENT.md"
    _write_text(p_deploy, deploy_md)

    # 07 Detailed code guide and Appendix
    dart_files = list_dart_files(LIB_DIR)
    ts_files = sorted([p for p in (REPO_ROOT / "functions" / "src").rglob("*.ts") if p.is_file()]) if (REPO_ROOT / "functions" / "src").exists() else []

    detailed_md = build_detailed_file_guide(dart_files, ts_files)
    p_detailed = DOCS_DIR / "07_DETAILED_CODE_GUIDE.md"
    _write_text(p_detailed, detailed_md)

    appendix_md = build_source_appendix(dart_files, ts_files)
    p_appendix = DOCS_DIR / "08_APPENDIX_SOURCE_CODE.md"
    _write_text(p_appendix, appendix_md)

    # Index + full doc
    index_md = (
        _md_h1("Documentation Index")
        + f"_Generated on {stamp}._\n\n"
        + _bullet(
            [
                "00_EXECUTIVE_SUMMARY.md",
                "01_TOOLS_TECH_STACK.md",
                "02_WORKFLOWS.md",
                "03_DATA_MODELS.md",
                "04_FIRESTORE_SCHEMA.md",
                "05_CODE_WALKTHROUGH.md",
                "06_DEPLOYMENT.md",
                "07_DETAILED_CODE_GUIDE.md",
                "08_APPENDIX_SOURCE_CODE.md",
            ]
        )
    )
    p_index = DOCS_DIR / "INDEX.md"
    _write_text(p_index, index_md)

    full_md = (
        exec_md
        + "\n---\n\n"
        + tools_md
        + "\n---\n\n"
        + workflows_md
        + "\n---\n\n"
        + models_md
        + "\n---\n\n"
        + schema_md
        + "\n---\n\n"
        + inv_md
        + "\n---\n\n"
        + deploy_md
    )
    p_full = DOCS_DIR / "FULL_DOCUMENTATION.md"
    _write_text(p_full, full_md)

    long_md = full_md + "\n---\n\n" + detailed_md + "\n---\n\n" + appendix_md
    p_long = DOCS_DIR / "FULL_DOCUMENTATION_LONG.md"
    _write_text(p_long, long_md)

    return {
        "executive_summary": p_exec,
        "tools_stack": p_tools,
        "workflows": p_workflows,
        "data_models": p_models,
        "firestore_schema": p_schema,
        "code_walkthrough": p_inv,
        "deployment": p_deploy,
        "detailed_guide": p_detailed,
        "appendix_source": p_appendix,
        "index": p_index,
        "full": p_full,
        "full_long": p_long,
    }


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0


def _add_two_column_slide(prs: Presentation, title: str, left: list[str], right: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title

    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(4.3), Inches(5.3))
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5.3))

    def fill(tb, items):
        tf = tb.text_frame
        tf.clear()
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0

    fill(left_box, left)
    fill(right_box, right)


def build_pptx(doc_paths: dict[str, Path]) -> Path:
    stamp = datetime.now().strftime("%Y-%m-%d")
    prs = Presentation()

    _add_title_slide(
        prs,
        "Cartiqo (SwiftCart)",
        f"Product + Technical Overview\nGenerated {stamp}",
    )

    _add_bullets_slide(
        prs,
        "Problem",
        [
            "Retail billing is slow and inconsistent when inventory, barcode labels, and billing rules are not centrally managed.",
            "Mall operations need a single dashboard to manage multiple managers, subscriptions, and sales reporting.",
        ],
    )

    _add_bullets_slide(
        prs,
        "Solution",
        [
            "Customer scan-to-bill app (Flutter) connected to a mall-specific product catalog in Firestore.",
            "Web portals for Admins and Mall Managers to manage operations, inventory, and analytics.",
            "Export-ready reporting (CSV) and print-ready artifacts (PDF labels/receipts).",
        ],
    )

    _add_two_column_slide(
        prs,
        "Modules (At a Glance)",
        left=[
            "Customer App",
            "Auth (Email/Phone OTP)",
            "Mall selection",
            "Barcode scanning + cart",
            "Checkout + payment records",
            "User profile + history",
        ],
        right=[
            "Web Portal",
            "Admin dashboard (malls/subscriptions/managers)",
            "Manager portal (inventory, barcodes, billing settings)",
            "Sales dashboard (filters, grouping, CSV export)",
            "Announcements + support requests",
        ],
    )

    # Tech stack slide from docs
    deps = load_pubspec_dependencies()
    core = [
        "Flutter (Dart)",
        "Firebase Auth + Firestore",
        "Firebase Hosting (web)",
        "Provider state management",
        "CSV + PDF/Printing utilities",
    ]
    extra = []
    for pkg in ["mobile_scanner", "qr_flutter", "csv", "pdf", "printing"]:
        if pkg in (deps.get("dependencies", {}) or {}):
            extra.append(pkg)
    _add_bullets_slide(
        prs,
        "Technology Stack",
        core + (["Key packages: " + ", ".join(extra)] if extra else []),
    )

    _add_bullets_slide(
        prs,
        "Personas",
        [
            "Customer/App User: scan items, build cart, checkout, view history",
            "Mall Manager: manage inventory, pricing, labels, promotions, sales analysis",
            "Admin: manage malls, subscriptions, manager accounts, support, and reporting",
        ],
    )

    _add_bullets_slide(
        prs,
        "Scan-to-Bill Workflow",
        [
            "Select mall (Firestore `malls`)",
            "Scan barcode (mobile_scanner)",
            "Map barcode to product (malls/{mallId}/barcodes + products)",
            "Build cart (Provider state)",
            "Checkout applies billing settings and writes payment records",
        ],
    )

    firestore = extract_firestore_collections()
    _add_bullets_slide(
        prs,
        "Firestore Data (High Level)",
        [
            "Top-level collections: " + ", ".join(firestore.get("collection", [])[:10]) + ("..." if len(firestore.get("collection", [])) > 10 else ""),
            "Malls subcollections: products, barcodes, payments, managers, promotions, staff_activity",
            "Users subcollections: bills, payments",
            "Support: announcements, support_requests",
        ],
    )

    fn_exports = extract_functions_exports()
    if fn_exports:
        _add_bullets_slide(
            prs,
            "Cloud Functions (OTP + Cleanup)",
            [
                f"Exports: {', '.join(fn_exports)}",
                "SMS OTP via Twilio (requires TWILIO_* env vars)",
                "Email OTP via Resend (requires RESEND_API_KEY)",
                "Scheduled cleanup jobs for OTP collections",
            ],
        )

    _add_bullets_slide(
        prs,
        "Security and Compliance (Talking Points)",
        [
            "Firestore is the source of truth; access is governed by security rules (should be reviewed per deployment).",
            "Sensitive keys must be stored as secrets/environment variables, not committed to the repo.",
            "Auditability: staff actions are logged to mall staff_activity (best effort).",
        ],
    )

    _add_bullets_slide(
        prs,
        "Reporting and Export",
        [
            "Manager sales dashboard: daily/weekly/monthly/yearly grouping and filters",
            "CSV export of dashboard summary + grouped analysis + underlying payments",
            "PDF/print services for labels and receipts",
        ],
    )

    _add_bullets_slide(
        prs,
        "Demo Flow (Suggested)",
        [
            "1) Open web portal: admin login and mall manager login panels",
            "2) Admin: add mall, set plan/subscription, create manager accounts",
            "3) Manager: login, add products, generate barcodes, view sales dashboard, export CSV",
            "4) App: login, select mall, scan products, checkout, view history",
        ],
    )

    _add_bullets_slide(
        prs,
        "Roadmap (Optional)",
        [
            "RBAC + Firestore security rules hardening",
            "Automated billing reconciliation and settlement reports",
            "Role-specific analytics dashboards and forecasting",
            "Multi-tenant auditing and alerting",
        ],
    )

    out = DOCS_DIR / "Cartiqo_Project_Presentation.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> None:
    doc_paths = build_docs()
    pptx_path = build_pptx(doc_paths)
    print("Generated docs:")
    for k, v in doc_paths.items():
        print(f"  {k}: {v}")
    print(f"Generated PPTX: {pptx_path}")


if __name__ == "__main__":
    main()
